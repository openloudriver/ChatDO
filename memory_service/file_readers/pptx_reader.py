"""
PPTX file reader using python-pptx.
"""
from pathlib import Path
from typing import Optional
from pptx import Presentation


def read_pptx(path: Path) -> Optional[str]:
    """
    Extract text from a PPTX file.
    
    Args:
        path: Path to the PPTX file
        
    Returns:
        Extracted text from all slides, or None if extraction fails
    """
    try:
        prs = Presentation(str(path))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if len(slide_text) > 1:  # More than just the slide number
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error reading PPTX {path}: {e}")
        return None

