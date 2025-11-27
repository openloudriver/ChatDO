"""
FastAPI server for ChatDO v1.5
Provides REST API and WebSocket endpoints for the ChatGPT-style UI
"""
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, WebSocket, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import sys
import json
import uuid
from uuid import uuid4
import re
import os
import logging
from pathlib import Path
from datetime import datetime, timedelta, timezone, date
from contextlib import asynccontextmanager
from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables from .env file in project root
env_path = Path(__file__).parent.parent / ".env"
load_dotenv(env_path)

# Add parent directory to path to import chatdo
sys.path.insert(0, str(Path(__file__).parent.parent))

from chatdo.config import load_target, TargetConfig
from chatdo.agents.main_agent import run_agent
from chatdo.memory.store import delete_thread_history, load_thread_history, save_thread_history, load_thread_sources, add_thread_source
from chatdo.executor import parse_tasks_block, apply_tasks
from server.uploads import handle_file_upload
from server.scraper import scrape_url
from server.ws import websocket_endpoint
from server.article_summary import extract_article
from chatdo.agents.main_agent import call_ai_router, ARTICLE_SUMMARY_SYSTEM_PROMPT, FILE_SUMMARY_SYSTEM_PROMPT
from server.services import impact_store, impact_templates_store
from server.services.impact_store import ImpactEntry
from server.services.impact_templates_store import ImpactTemplate, get_template_file_path
from server.services.template_engine import template_store, extract_template_text, analyze_template, autofill_template
from server.services.template_engine import template_store, extract_template_text, analyze_template, autofill_template

# Retention settings
RETENTION_DAYS = int(os.getenv("CHATDO_TRASH_RETENTION_DAYS", "30"))

# Task execution constants
TASKS_START = "<TASKS>"
TASKS_END = "</TASKS>"


def split_tasks_block(text: str) -> tuple[str, Optional[str]]:
    """
    Split a ChatDO response into (human_text, tasks_json_str | None).
    
    - human_text: original text with any <TASKS>...</TASKS> block removed.
    - tasks_json_str: the raw JSON string inside the <TASKS> block, or None.
    """
    start = text.find(TASKS_START)
    end = text.find(TASKS_END)
    
    if start == -1 or end == -1:
        return text, None
    
    json_str = text[start + len(TASKS_START) : end].strip()
    
    # Remove the block from the human-facing text
    human = (text[:start] + text[end + len(TASKS_END) :]).strip()
    
    return human, json_str


# Chat metadata models
class Chat(BaseModel):
    id: str
    project_id: str
    title: str
    thread_id: str
    created_at: str
    updated_at: str
    trashed: bool = False
    trashed_at: Optional[str] = None


class ChatCreate(BaseModel):
    project_id: str
    title: str
    thread_id: str


class ChatUpdate(BaseModel):
    title: Optional[str] = None


# Helper functions for chat persistence
def get_chats_path() -> Path:
    """Get the path to chats.json"""
    return Path(__file__).parent / "data" / "chats.json"


def load_chats() -> List[dict]:
    """Load chats from chats.json"""
    chats_path = get_chats_path()
    if not chats_path.exists():
        return []
    
    with open(chats_path, "r") as f:
        return json.load(f)


def save_chats(chats: List[dict]) -> None:
    """Save chats to chats.json"""
    chats_path = get_chats_path()
    chats_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(chats_path, "w") as f:
        json.dump(chats, f, indent=2)


def get_active_chats(chats: List[dict]) -> List[dict]:
    """Get chats that are not trashed"""
    return [c for c in chats if not c.get("trashed", False)]


def get_trashed_chats(chats: List[dict]) -> List[dict]:
    """Get chats that are trashed"""
    return [c for c in chats if c.get("trashed", False)]


def now_iso() -> str:
    """Get current UTC time as ISO8601 string"""
    return datetime.now(timezone.utc).isoformat()


def purge_trashed_chats() -> int:
    """
    Permanently delete chats that have been in trash longer than RETENTION_DAYS.
    Returns the number of chats purged.
    """
    chats = load_chats()
    cutoff_date = datetime.now(timezone.utc) - timedelta(days=RETENTION_DAYS)
    purged_count = 0
    
    chats_to_purge = []
    for chat in chats:
        if chat.get("trashed", False) and chat.get("trashed_at"):
            try:
                trashed_at = datetime.fromisoformat(chat["trashed_at"].replace("Z", "+00:00"))
                if trashed_at < cutoff_date:
                    chats_to_purge.append(chat)
            except (ValueError, KeyError):
                # Invalid date format, skip
                continue
    
    # Purge each chat
    for chat in chats_to_purge:
        # Delete thread history
        project = None
        projects = load_projects()
        for p in projects:
            if p.get("id") == chat.get("project_id"):
                project = p
                break
        
        if project:
            target_name = project.get("default_target", "general")
            thread_id = chat.get("thread_id")
            if thread_id:
                try:
                    delete_thread_history(target_name, thread_id)
                except Exception as e:
                    print(f"Warning: Failed to delete thread history for {thread_id}: {e}")
        
        # Remove from chats list
        chats = [c for c in chats if c.get("id") != chat.get("id")]
        purged_count += 1
    
    if purged_count > 0:
        save_chats(chats)
    
    return purged_count


@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup: purge old trashed chats
    purged = purge_trashed_chats()
    if purged > 0:
        print(f"Purged {purged} trashed chats older than {RETENTION_DAYS} days")
    yield
    # Shutdown (if needed)


app = FastAPI(title="ChatDO API", version="1.5.0", lifespan=lifespan)

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000", "http://127.0.0.1:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class ChatRequest(BaseModel):
    rag_file_ids: Optional[List[str]] = None  # IDs of RAG files to use as context
    project_id: str
    conversation_id: str
    target_name: str
    message: str
    force_search: bool = False  # Explicit "Search web" trigger


class ChatResponse(BaseModel):
    reply: str
    model_used: Optional[str] = None
    provider: Optional[str] = None
    message_type: Optional[str] = None  # 'text' (default), 'web_search_results', or 'article_card'
    message_data: Optional[Dict[str, Any]] = None  # Structured data for special message types
    sources: Optional[List[str]] = None  # Source tags (e.g., ["Memory-DRR", "Brave Search"])


class ArticleSummaryRequest(BaseModel):
    url: str
    title: Optional[str] = None  # Optional title from UI
    conversation_id: Optional[str] = None  # Optional conversation_id to save to memory store
    project_id: Optional[str] = None  # Optional project_id to determine target_name
    force_mode: Optional[str] = None  # Deprecated - kept for backward compatibility but not used
    privacy_mode: bool = False  # Deprecated - kept for backward compatibility but ignored


class ArticleSummaryResponse(BaseModel):
    message_type: str = "article_card"
    message_data: Dict[str, Any]
    model: str = "Trafilatura + GPT-5"
    provider: str = "trafilatura-gpt5"
    model_label: Optional[str] = None  # Human-readable model label (e.g., "YouTube transcript + GPT-5")


# RAG File model
class RagFile(BaseModel):
    id: str
    chat_id: str  # conversation_id
    filename: str
    mime_type: str
    size: int
    created_at: str
    text_path: Optional[str] = None  # Path to extracted text file
    path: Optional[str] = None  # Path to original file (for preview)
    text_extracted: bool = False
    error: Optional[str] = None


class NewConversationRequest(BaseModel):
    project_id: str


class NewConversationResponse(BaseModel):
    conversation_id: str


class Project(BaseModel):
    id: str
    name: str
    default_target: str
    sort_index: int = 0


class ProjectCreate(BaseModel):
    name: str


class ProjectUpdate(BaseModel):
    name: str


# Helper functions for project persistence
def get_projects_path() -> Path:
    """Get the path to projects.json"""
    return Path(__file__).parent / "data" / "projects.json"


def load_projects() -> List[dict]:
    """Load projects from projects.json, ensuring sort_index and memory_sources exist"""
    from server.services import projects_config
    return projects_config.load_projects()


def save_projects(projects: List[dict]) -> None:
    """Save projects to projects.json"""
    from server.services import projects_config
    projects_config.save_projects(projects)


def generate_slug(name: str) -> str:
    """Generate a URL-friendly slug from a name"""
    # Convert to lowercase and replace spaces/special chars with hyphens
    slug = re.sub(r'[^\w\s-]', '', name.lower())
    slug = re.sub(r'[-\s]+', '-', slug)
    return slug.strip('-')


def generate_project_id() -> str:
    """Generate a unique project ID"""
    return str(uuid.uuid4())


@app.get("/")
async def root():
    return {"message": "ChatDO API v1.5", "status": "running"}


@app.get("/api/projects", response_model=List[Project])
async def get_projects():
    """Return list of available projects, ensuring General exists"""
    projects = load_projects()
    
    # Ensure General project exists
    general_project = next((p for p in projects if p.get("name") == "General"), None)
    if not general_project:
        # Create General project
        general_project = {
            "id": str(uuid.uuid4()),
            "name": "General",
            "default_target": "general",
            "sort_index": 0,
            "created_at": now_iso(),
            "updated_at": now_iso()
        }
        projects.insert(0, general_project)
        save_projects(projects)
    
    return projects


@app.post("/api/projects", response_model=Project)
async def create_project(project_data: ProjectCreate):
    """Create a new project"""
    projects = load_projects()
    
    # Generate ID and slug
    project_id = generate_project_id()
    slug = generate_slug(project_data.name)
    
    # Create new project with sort_index at the end
    new_project = {
        "id": project_id,
        "name": project_data.name.strip(),
        "default_target": "general",  # Default target for new projects
        "sort_index": len(projects),  # Add at the end
        "memory_sources": []  # New projects start with no memory sources
    }
    
    projects.append(new_project)
    save_projects(projects)
    
    return new_project


@app.patch("/api/projects/{project_id}", response_model=Project)
async def update_project(project_id: str, project_data: ProjectUpdate):
    """Update an existing project"""
    projects = load_projects()
    
    # Find project by ID
    project_index = None
    for i, p in enumerate(projects):
        if p.get("id") == project_id:
            project_index = i
            break
    
    if project_index is None:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Update project name
    projects[project_index]["name"] = project_data.name.strip()
    save_projects(projects)
    
    return projects[project_index]


@app.delete("/api/projects/{project_id}")
async def delete_project(project_id: str):
    """Delete a project"""
    projects = load_projects()
    
    # Find and remove project
    original_count = len(projects)
    projects = [p for p in projects if p.get("id") != project_id]
    
    if len(projects) == original_count:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Reindex remaining projects to maintain sequential sort_index
    for i, project in enumerate(projects):
        project["sort_index"] = i
    
    save_projects(projects)
    return {"success": True}


@app.post("/api/projects/reorder", response_model=List[Project])
async def reorder_projects(order: List[str]):
    """
    Reorder projects based on the provided array of project IDs.
    Updates sort_index for each project to match the new order.
    """
    projects = load_projects()
    
    # Create a map of project ID to project for quick lookup
    project_map = {p.get("id"): p for p in projects}
    
    # Update sort_index for projects in the order list
    reordered_projects = []
    for index, project_id in enumerate(order):
        if project_id in project_map:
            project = project_map[project_id]
            project["sort_index"] = index
            reordered_projects.append(project)
    
    # Add any projects not in the order list (defensive - shouldn't happen)
    for project in projects:
        if project.get("id") not in order:
            project["sort_index"] = len(reordered_projects)
            reordered_projects.append(project)
    
    # Sort to ensure consistency
    reordered_projects.sort(key=lambda p: (p.get("sort_index", 0), p.get("name", "")))
    
    save_projects(reordered_projects)
    return reordered_projects


@app.post("/api/new_conversation", response_model=NewConversationResponse)
async def new_conversation(request: NewConversationRequest):
    """Create a new conversation thread"""
    conversation_id = str(uuid.uuid4())
    
    # Create chat metadata entry
    chats = load_chats()
    projects = load_projects()
    project = next((p for p in projects if p.get("id") == request.project_id), None)
    target_name = project.get("default_target", "general") if project else "general"
    
    new_chat = {
        "id": conversation_id,
        "project_id": request.project_id,
        "title": "New Chat",
        "thread_id": conversation_id,
        "created_at": now_iso(),
        "updated_at": now_iso(),
        "trashed": False,
        "trashed_at": None
    }
    
    chats.append(new_chat)
    save_chats(chats)
    
    return NewConversationResponse(conversation_id=conversation_id)


def extract_urls(text: str) -> List[str]:
    """Extract all URLs from text."""
    url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
    urls = re.findall(url_pattern, text)
    return urls


@app.post("/api/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    """
    Main chat endpoint - hybrid approach:
    - URL summaries: unchanged (YouTube transcript, yt-dlp+Whisper, Trafilatura)
    - Explicit search commands: return Brave "Top Results" card
    - Normal chat: GPT-5 with smart auto-search (Brave in background when needed)
    """
    try:
        # 1. Check for single URL with "summarize" keyword - route to article summary (unchanged)
        urls = extract_urls(request.message)
        if len(urls) == 1 and ("summarize" in request.message.lower() or "summary" in request.message.lower()):
            # Get project_id from conversation if not in request
            project_id = getattr(request, 'project_id', None)
            if not project_id and request.conversation_id:
                chats = load_chats()
                chat = next((c for c in chats if c.get("id") == request.conversation_id), None)
                if chat:
                    project_id = chat.get("project_id")
            
            article_request = ArticleSummaryRequest(
                url=urls[0],
                conversation_id=request.conversation_id,
                project_id=project_id
            )
            result = await summarize_article(article_request)
            return ChatResponse(
                reply=result.message_data.get("summary", ""),
                message_type=result.message_type,
                message_data=result.message_data,
                model_used=result.model,
                provider=result.provider
            )
        
        # 2. Check for explicit search (force_search flag or search: prefix)
        message_lower = request.message.lower().strip()
        explicit_search_patterns = ["search:", "web search:", "brave:", "find:"]
        is_explicit_search = request.force_search or any(message_lower.startswith(pattern) for pattern in explicit_search_patterns)
        
        if is_explicit_search:
            # Extract query from command
            query = request.message
            for pattern in explicit_search_patterns:
                if message_lower.startswith(pattern):
                    query = request.message[len(pattern):].strip()
                    break
            
            if not query:
                query = request.message
            
            # Perform explicit Brave search and return Top Results card
            from chatdo.tools import web_search
            try:
                search_results = web_search.search_web(query, max_results=10)
                if search_results and len(search_results) > 0:
                    structured_result = {
                        "type": "web_search_results",
                        "query": query,
                        "provider": "brave",
                        "results": search_results
                    }
                    
                    # Save to memory store
                    if request.conversation_id:
                        target_cfg = load_target(request.target_name)
                        history = load_thread_history(target_cfg.name, request.conversation_id)
                        history.append({"role": "user", "content": request.message})
                        history.append({
                            "role": "assistant",
                            "content": "",
                            "type": "web_search_results",
                            "data": structured_result,
                            "model": "Brave Search",
                            "provider": "brave_search"
                        })
                        save_thread_history(target_cfg.name, request.conversation_id, history)
                    
                    return ChatResponse(
                        reply="",
                        message_type="web_search_results",
                        message_data=structured_result,
                        model_used="Brave Search",
                        provider="brave_search",
                        sources=["Brave Search"]
                    )
                else:
                    return ChatResponse(
                        reply="No search results found. Please try a different query.",
                        model_used="Brave Search",
                        provider="brave_search"
                    )
            except Exception as e:
                error_msg = f"Web search failed: {str(e)}"
                return ChatResponse(
                    reply=error_msg,
                    model_used="Brave Search",
                    provider="brave_search"
                )
        
        # 3. Normal chat with smart auto-search (or RAG if provided)
        # Build RAG context if RAG files are provided
        rag_context = ""
        has_rag_context = False
        if request.rag_file_ids:
            print(f"[RAG] Received {len(request.rag_file_ids)} RAG file IDs: {request.rag_file_ids}")
            rag_context = build_rag_context(request.rag_file_ids, request.message, chat_id=request.conversation_id)
            if rag_context:
                has_rag_context = True
                print(f"[RAG] Context built successfully, length: {len(rag_context)}")
            else:
                print(f"[RAG] Warning: No context was built despite {len(request.rag_file_ids)} file IDs provided")
        
        # Load target configuration
        target_cfg = load_target(request.target_name)
        
        # If RAG context is provided, use run_agent (existing RAG flow)
        # Otherwise, use smart chat with auto-search
        if has_rag_context:
            # Prepend RAG context to user message
            user_message = f"{rag_context}\n\nUser question: {request.message}"
            
            # Use run_agent for RAG (existing flow)
            raw_result, model_display, provider = run_agent(
                target=target_cfg,
                task=user_message,  # This includes RAG context
                thread_id=request.conversation_id if request.conversation_id else None,
                skip_web_search=True  # Skip web search when RAG is available
            )
        else:
            # Use smart chat with auto-search for normal chat
            from server.services.chat_with_smart_search import chat_with_smart_search
            
            # Get project_id from conversation if available
            project_id = request.project_id if hasattr(request, 'project_id') else None
            if not project_id and request.conversation_id:
                chats = load_chats()
                chat = next((c for c in chats if c.get("id") == request.conversation_id), None)
                if chat:
                    project_id = chat.get("project_id")
            
            result = await chat_with_smart_search(
                user_message=request.message,
                target_name=target_cfg.name,
                thread_id=request.conversation_id if request.conversation_id else None,
                project_id=project_id
            )
            
            # If result has tasks or needs special handling, route to run_agent
            # Otherwise return the smart chat result
            if result.get("type") == "assistant_message":
                # Check if content contains <TASKS> block (needs task execution)
                content = result.get("content", "")
                human_text, tasks_json = split_tasks_block(content)
                
                if tasks_json:
                    # Has tasks - need to execute them
                    # Continue with task execution flow below
                    raw_result = content
                    model_display = result.get("model", "GPT-5")
                    provider = result.get("provider", "openai-gpt5")
                else:
                    # No tasks - return smart chat result directly
                    # Add meta information for frontend
                    return ChatResponse(
                        reply=human_text,  # Use human_text (content without tasks)
                        message_type="text",
                        message_data={
                            "content": human_text,
                            "meta": result.get("meta", {})
                        },
                        model_used=result.get("model", "GPT-5"),
                        provider=result.get("provider", "openai-gpt5"),
                        sources=result.get("sources")
                    )
            else:
                # Unexpected result type, fall back to run_agent
                raw_result, model_display, provider = run_agent(
                    target=target_cfg,
                    task=request.message,
                    thread_id=request.conversation_id if request.conversation_id else None,
                    skip_web_search=False
                )
        
        # Continue with existing flow for RAG or tasks
        # NOTE: run_agent now filters out RAG context automatically, so we don't need to fix it here
        # But we'll keep this as a safety net in case any RAG context slips through
        if has_rag_context and request.conversation_id:
            print(f"[DIAG] REST: RAG context was used, but run_agent should have filtered it. Verifying...")
            try:
                target_name_save = target_cfg.name
                thread_id = request.conversation_id
                history = load_thread_history(target_name_save, thread_id)
                
                # Check if last user message still has RAG context
                if len(history) >= 2:
                    user_idx = len(history) - 2
                    if history[user_idx].get("role") == "user":
                        user_content = history[user_idx].get("content", "")
                        # Check if it contains RAG context markers
                        if "You have access to the following reference documents" in user_content or "----\nSource:" in user_content:
                            print(f"[DIAG] REST: WARNING - User message still contains RAG context! Fixing now...")
                            history[user_idx]["content"] = request.message
                            save_thread_history(target_name_save, thread_id, history)
                            print(f"[DIAG] REST: ✅ FIXED user message (removed RAG context)")
            except Exception as e:
                print(f"[DIAG] REST: Error verifying/fixing user message: {e}")
        
        # 2) Check if result is structured (web_search_results or article_card)
        if isinstance(raw_result, dict):
            result_type = raw_result.get("type")
            if result_type == "web_search_results":
                # Return structured web search results
                return ChatResponse(
                    reply="",  # Empty reply for structured messages
                    message_type="web_search_results",
                    message_data=raw_result,
                    model_used=model_display,
                    provider=provider
                )
            elif result_type == "article_card":
                # Return structured article card results
                return ChatResponse(
                    reply="",  # Empty reply for structured messages
                    message_type="article_card",
                    message_data=raw_result,
                    model_used=model_display,
                    provider=provider
                )
        
        # 3) Split out any <TASKS> block
        human_text, tasks_json = split_tasks_block(raw_result)
        
        # 3.5) Guardrail: Check for missing citations in RAG responses
        if has_rag_context and not tasks_json and human_text:
            if check_missing_citations(human_text):
                print(f"[RAG] Guardrail: Missing citations detected, requesting fix-up...")
                human_text = await fix_missing_citations(human_text, rag_context, model_display, provider)
                # Re-check after fix-up
                if check_missing_citations(human_text):
                    print(f"[RAG] Guardrail: Still missing citations after fix-up, but proceeding with response")
                else:
                    print(f"[RAG] Guardrail: ✅ Citations added successfully")
        
        # 4) If RAG context was used, return as structured rag_response
        if has_rag_context and not tasks_json:
            # Extract source file names from RAG context
            # IMPORTANT: Return sources in the SAME ORDER as rag_file_ids to match RAG tray numbering
            rag_files = load_rag_files(request.conversation_id) if request.rag_file_ids else []
            # Create a lookup dict for fast access
            rag_files_by_id = {f.get("id"): f for f in rag_files}
            # Return sources in the order of rag_file_ids (matches RAG tray order)
            source_files = [rag_files_by_id.get(fid, {}).get("filename", "") 
                          for fid in (request.rag_file_ids or [])
                          if fid in rag_files_by_id and rag_files_by_id[fid].get("text_extracted")]
            
            # Log raw RAG output to verify markdown headings are present
            print(f"[RAG] Raw RAG output (first 500 chars):\n{human_text[:500]}")
            if "###" in human_text:
                print("[RAG] ✅ Markdown headings (###) detected in response")
            else:
                print("[RAG] ⚠️  WARNING: No markdown headings (###) found in response")
            
            # Update the last message in memory store to have structured type
            # (run_agent already saved it as a regular message, so we update it)
            if request.conversation_id and request.project_id:
                try:
                    projects = load_projects()
                    project = next((p for p in projects if p.get("id") == request.project_id), None)
                    if project:
                        target_name = project.get("default_target", "general")
                        thread_id = request.conversation_id
                        
                        history = load_thread_history(target_name, thread_id)
                        # Find the last assistant message (saved by run_agent) and update it with structured type
                        updated = False
                        for i in range(len(history) - 1, -1, -1):
                            if history[i].get("role") == "assistant":
                                # Update existing message with structured type
                                history[i]["type"] = "rag_response"
                                history[i]["data"] = {
                                    "content": human_text,
                                    "sources": source_files
                                }
                                history[i]["model"] = model_display
                                history[i]["provider"] = provider
                                updated = True
                                print(f"[RAG] Updated message at index {i} with type=rag_response, data keys: {list(history[i].get('data', {}).keys())}")
                                break
                        if updated:
                            save_thread_history(target_name, thread_id, history)
                            print(f"[RAG] Saved updated history with {len(history)} messages")
                        else:
                            print(f"[RAG] Warning: Could not find assistant message to update in history of {len(history)} messages")
                except Exception as e:
                    print(f"Warning: Failed to update RAG response in memory store: {e}")
            
            return ChatResponse(
                reply=human_text,
                message_type="rag_response",
                message_data={
                    "content": human_text,
                    "sources": source_files
                },
                model_used=model_display,
                provider=provider,
                sources=["RAG-Upload"] if source_files else None
            )
        
        # 5) If no tasks, return (run_agent already saved the message)
        if not tasks_json:
            return ChatResponse(reply=human_text, model_used=model_display, provider=provider)
        
        # 6) Parse tasks from JSON
        try:
            tasks = parse_tasks_block(tasks_json)
        except Exception as e:
            # If parsing fails, show error but still return the human text
            error_note = f"\n\n---\nExecutor error: could not parse tasks JSON ({e})."
            return ChatResponse(reply=human_text + error_note)
        
        # 7) Execute tasks against the target repo
        exec_result = apply_tasks(target_cfg, tasks)
        
        # 8) Build a human-readable summary
        summary_lines = [exec_result.summary()]
        for r in exec_result.results:
            prefix = "✅" if r.status == "success" else "❌"
            summary_lines.append(f"{prefix} {r.message}")
        
        summary_text = "\n".join(summary_lines)
        
        # 8) Append summary to what the user sees
        final_message = human_text + "\n\n---\nExecutor summary:\n" + summary_text
        
        return ChatResponse(reply=final_message, model_used=model_display, provider=provider)
    
    except Exception as e:
        import traceback
        error_detail = str(e)
        traceback.print_exc()  # Print full traceback to server logs
        raise HTTPException(status_code=500, detail=error_detail)


@app.post("/api/article/summary", response_model=ArticleSummaryResponse)
async def summarize_article(request: ArticleSummaryRequest):
    """
    Summarize a URL (web page or video) using deterministic 2-tier routing.
    
    Routing Rules:
    1. If URL is YouTube (youtube.com or youtu.be)
       → Tier 1: YouTube transcript API + GPT-5
       → NO FALLBACK: If transcript API fails, return error
    
    2. Else if URL is a video host (rumble.com, bitchute.com, archive.org)
       → Tier 2: yt-dlp + Whisper-small-FP16 + GPT-5 (M1-optimized)
       → NO FALLBACK: If audio pipeline fails, return error
    
    3. Else
       → Web page: Trafilatura + GPT-5
    
    Model labels:
    - "YouTube transcript + GPT-5" (Tier 1)
    - "yt-dlp + Whisper-small-FP16 + GPT-5" (Tier 2, M1-optimized)
    - "Trafilatura + GPT-5" (Web pages)
    
    Returns a structured article_card message.
    """
    try:
        # Validate URL
        if not request.url.startswith(("http://", "https://")):
            raise HTTPException(status_code=400, detail="Invalid URL. Must start with http:// or https://")
        
        # Deterministic URL classification
        from server.utils.url_classification import (
            is_youtube_url,
            is_other_video_host,
        )
        from server.services.youtube_transcript import (
            get_youtube_transcript,
            YouTubeTranscriptError,
        )
        from server.services.video_transcription import get_transcript_from_url
        from urllib.parse import urlparse
        
        # Get source text based on URL type
        article_text = None
        article_title = None
        article_site_name = None
        article_published = None  # Store published date from extraction
        model_label = None
        
        if is_youtube_url(request.url):
            # --- Tier 1: YouTube-only ---
            # youtube-transcript-api + GPT-5
            # NO FALLBACK: If transcript API fails, return error (do NOT call yt-dlp or Whisper)
            try:
                article_text = get_youtube_transcript(request.url)
                parsed = urlparse(request.url)
                article_site_name = "Video"
                article_title = f"Video from {parsed.netloc}"
                model_label = "YouTube transcript + GPT-5"
            except YouTubeTranscriptError as e:
                # Explicitly do NOT fall back to audio pipeline
                logger.warning("Tier 1 (YouTube transcript) failed: url=%s error=%s", request.url, e)
                raise HTTPException(
                    status_code=400,
                    detail=f"Could not obtain YouTube transcript: {str(e)}",
                ) from e
            except Exception as e:
                logger.exception("Tier 1 (YouTube transcript) unexpected error: url=%s", request.url)
                raise HTTPException(
                    status_code=500,
                    detail=f"Error retrieving YouTube transcript: {str(e)}",
                ) from e
        elif is_other_video_host(request.url):
            # --- Tier 2: Non-YouTube video hosts ---
            # Rumble / BitChute / Archive.org
            # yt-dlp → Whisper-small-FP16 (M1-optimized) → GPT-5
            # NO FALLBACK: If audio pipeline fails, return error (do NOT try web extraction)
            try:
                article_text = await get_transcript_from_url(request.url)
                parsed = urlparse(request.url)
                article_site_name = "Video"
                article_title = f"Video from {parsed.netloc}"
                model_label = "yt-dlp + Whisper-small-FP16 + GPT-5"
            except ValueError as ve:
                raise HTTPException(status_code=400, detail=str(ve))
            except RuntimeError as re:
                logger.exception("Tier 2 (audio pipeline) failed: url=%s", request.url)
                raise HTTPException(
                    status_code=502,
                    detail="Couldn't obtain transcript for this video. The audio download or transcription failed.",
                ) from re
            except Exception as e:
                raise HTTPException(
                    status_code=500,
                    detail="Unexpected error while summarizing this video URL.",
                ) from e
        else:
            # --- Web page ---
            # Trafilatura + GPT-5
            logging.info(f"[Summary] Processing as web page: {request.url}")
            article_data = extract_article(request.url)
            
            if article_data.get("error"):
                raise HTTPException(status_code=400, detail=article_data["error"])
            
            if not article_data.get("text"):
                raise HTTPException(status_code=400, detail="Could not extract article content.")
            
            article_text = article_data["text"]
            article_title = article_data.get("title")
            article_site_name = article_data.get("site_name")
            article_published = article_data.get("published")  # Extract published date
            model_label = "Trafilatura + GPT-5"
        
        if not article_text:
            raise HTTPException(status_code=502, detail="Content could not be processed.")
        
        # Truncate text to safe length (first 10k chars)
        article_text = article_text[:10000]
        
        # Summarize with GPT-5 (always use GPT-5, no privacy mode)
        is_video = is_youtube_url(request.url) or is_other_video_host(request.url)
        source_type = "video transcript" if is_video else "article"
        user_prompt = f"""Please summarize the following {source_type}:

{article_text}

Provide:
1. A 2–4 sentence summary paragraph
2. 3–5 key bullet points
3. 1–2 sentences on why this matters or its significance

Format your response clearly with the summary first, then key points (as bullet points), then "Why This Matters:" followed by your analysis.

Keep it concise, neutral, and factual."""
        
        # Call AI Router with summarize_article intent
        messages = [
            {"role": "system", "content": ARTICLE_SUMMARY_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt}
        ]
        
        ai_router_url = os.getenv("AI_ROUTER_URL", "http://localhost:8081/v1/ai/run")
        payload = {
            "role": "chatdo",
            "intent": "summarize_article",
            "priority": "high",
            "privacyLevel": "normal",
            "costTier": "standard",
            "input": {
                "messages": messages,
            },
        }
        
        import requests
        resp = requests.post(ai_router_url, json=payload, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        
        if not data.get("ok"):
            raise HTTPException(status_code=500, detail=f"AI Router error: {data.get('error')}")
        
        assistant_messages = data["output"]["messages"]
        if not assistant_messages or len(assistant_messages) == 0:
            raise HTTPException(status_code=500, detail="No response from AI Router")
        
        summary_text = assistant_messages[0].get("content", "")
        
        # Parse summary into summary paragraph and bullet points
        # Simple parsing: look for bullet points (lines starting with - or •)
        lines = summary_text.split("\n")
        summary_paragraph = ""
        key_points = []
        why_matters = ""
        
        current_section = "summary"
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            line_lower = line.lower()
            
            # Detect "why this matters" section - be more flexible
            if ("why" in line_lower and ("matter" in line_lower or "context" in line_lower or "significance" in line_lower)):
                current_section = "why_matters"
                # If the line has content after a colon, extract it
                if ":" in line:
                    content = line.split(":", 1)[1].strip()
                    if content:
                        why_matters = content
                continue
            
            # Detect bullet points
            if line.startswith("-") or line.startswith("•") or (line.startswith("*") and len(line) > 1):
                bullet_text = line.lstrip("-•* ").strip()
                if bullet_text:
                    if current_section == "why_matters":
                        # Bullet point in why_matters section
                        if why_matters:
                            why_matters += " " + bullet_text
                        else:
                            why_matters = bullet_text
                    else:
                        key_points.append(bullet_text)
            elif current_section == "summary":
                if not summary_paragraph:
                    summary_paragraph = line
                else:
                    summary_paragraph += " " + line
            elif current_section == "why_matters":
                if why_matters:
                    why_matters += " " + line
                else:
                    why_matters = line
        
        # If no structured parsing worked, use the whole text as summary
        if not summary_paragraph:
            summary_paragraph = summary_text[:500]  # First 500 chars as summary
        
        # Get title with proper fallback
        title = request.title or article_title
        if not title or title.strip() == "":
            from urllib.parse import urlparse
            domain = urlparse(request.url).netloc.replace("www.", "")
            title = f"{'Video' if is_video else 'Article'} from {domain}"
        
        # Build response
        message_data = {
            "url": request.url,
            "title": title,
            "siteName": article_site_name or "",
            "published": article_published,  # Use extracted date for web pages, None for videos
            "summary": summary_paragraph,
            "keyPoints": key_points if key_points else [],
            "whyMatters": why_matters if why_matters else None,
        }
        
        # Save to memory store if conversation_id is provided
        if request.conversation_id and request.project_id:
            try:
                projects = load_projects()
                project = next((p for p in projects if p.get("id") == request.project_id), None)
                if project:
                    target_name = project.get("default_target", "general")
                    thread_id = request.conversation_id
                    
                    # Load existing history
                    history = load_thread_history(target_name, thread_id)
                    
                    # Add user message (if URL was provided by user)
                    user_message = f"Summarize: {request.url}"
                    history.append({"role": "user", "content": user_message})
                    
                    # Add assistant message with structured data
                    # Determine provider for history
                    is_video_for_history = is_youtube_url(request.url) or is_other_video_host(request.url)
                    if is_video_for_history:
                        if is_youtube_url(request.url):
                            provider_for_history = "youtube-transcript-api"
                        else:
                            provider_for_history = "yt-dlp-whisper-small"
                    else:
                        provider_for_history = "trafilatura-gpt5"
                    
                    assistant_message = {
                        "role": "assistant",
                        "content": "",  # Empty content for structured messages
                        "type": "article_card",
                        "data": message_data,
                        "model": model_label or ("Trafilatura + GPT-5" if not is_video_for_history else ("YouTube transcript + GPT-5" if is_youtube_url(request.url) else "yt-dlp + Whisper-small-FP16 + GPT-5")),
                        "provider": provider_for_history
                    }
                    history.append(assistant_message)
                    
                    # Save back to memory store
                    save_thread_history(target_name, thread_id, history)
                    
                    # Also save as a source
                    import uuid as uuid_lib
                    source = {
                        "id": str(uuid_lib.uuid4()),
                        "kind": "url",
                        "title": message_data.get("title") or title,
                        "description": summary_paragraph[:200] if summary_paragraph else None,
                        "url": request.url,
                        "createdAt": datetime.now(timezone.utc).isoformat(),
                        "meta": {
                            "siteName": article_site_name,
                            "published": article_published,  # Use extracted date for web pages, None for videos
                        }
                    }
                    add_thread_source(target_name, thread_id, source)
            except Exception as e:
                # Don't fail the request if saving to memory store fails
                print(f"Warning: Failed to save article summary to memory store: {e}")
        
        # Determine model/provider name based on mode
        # model_label should already be set correctly by the routing logic above
        is_video = is_youtube_url(request.url) or is_other_video_host(request.url)
        model_name = model_label or ("Trafilatura + GPT-5" if not is_video else ("YouTube transcript + GPT-5" if is_youtube_url(request.url) else "yt-dlp + Whisper-small-FP16 + GPT-5"))
        if is_video:
            if is_youtube_url(request.url):
                provider_name = "youtube-transcript-api"
            else:
                provider_name = "yt-dlp-whisper-small"
        else:
            provider_name = "trafilatura-gpt5"
        
        return ArticleSummaryResponse(
            message_data=message_data,
            model=model_name,
            provider=provider_name,
            model_label=model_label
        )
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error summarizing article: {str(e)}")


# Removed auto-summarize endpoint - files are now used for RAG context instead


@app.post("/api/upload")
async def upload_file(
    project_id: str = Form(...),
    conversation_id: str = Form(...),
    file: UploadFile = File(...)
):
    """
    Handle file uploads (legacy endpoint for one-off file sharing in conversations)
    Saves to uploads/<project_id>/<conversation_id>/<uuid>.<ext>
    Extracts text if PDF/Word/Excel/PPT/image
    Note: For RAG context files, use /api/rag/files instead
    """
    try:
        result = await handle_file_upload(project_id, conversation_id, file)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# RAG File Storage
RAG_FILES_DIR = Path(__file__).parent.parent / "rag_files"
RAG_FILES_DIR.mkdir(exist_ok=True)

def get_rag_files_path(chat_id: str) -> Path:
    """Get path to RAG files JSON for a chat"""
    return RAG_FILES_DIR / f"{chat_id}.json"

def load_rag_files(chat_id: str) -> List[Dict[str, Any]]:
    """Load RAG files for a chat"""
    rag_path = get_rag_files_path(chat_id)
    if not rag_path.exists():
        return []
    try:
        with open(rag_path, 'r') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading RAG files: {e}")
        return []

def save_rag_files(chat_id: str, files: List[Dict[str, Any]]):
    """Save RAG files for a chat"""
    rag_path = get_rag_files_path(chat_id)
    with open(rag_path, 'w') as f:
        json.dump(files, f, indent=2)

def chunk_text(text: str, chunk_size: int = 2000) -> List[str]:
    """Simple text chunking for RAG"""
    chunks = []
    words = text.split()
    current_chunk = []
    current_length = 0
    
    for word in words:
        word_length = len(word) + 1  # +1 for space
        if current_length + word_length > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_length = word_length
        else:
            current_chunk.append(word)
            current_length += word_length
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def build_rag_context(rag_file_ids: List[str], user_message: str, chat_id: Optional[str] = None) -> str:
    """Build RAG context message from file IDs"""
    if not rag_file_ids:
        return ""
    
    print(f"[RAG] Building context for {len(rag_file_ids)} files, chat_id={chat_id}")
    
    context_parts = [
        "You have access to the following reference documents for this conversation. Use them as primary sources when answering.\n",
        "\n",
        "IMPORTANT: When answering questions about these documents, format your response using markdown with level-3 headings (###) for each major section.\n",
        "Use this structure where applicable:\n",
        "- ### Summary\n",
        "- ### Highlights\n",
        "- ### Analysis (or similar section titles as needed)\n",
        "- ### Recommendation\n",
        "- ### Bottom line\n",
        "\n",
        "Under each heading, use bullet points where helpful. The exact section titles can vary based on the question, but always use markdown headings (### Heading) for every section title instead of plain text lines.\n",
        "\n",
        "FORMATTING RULES (MANDATORY):\n",
        "\n",
        "- Use markdown headings and bullet lists.\n",
        "- For EVERY bullet or sentence that uses information from the sources, you MUST include one or more inline citations in square brackets, like `[1]` or `[1, 3]`, at the **end** of the sentence or bullet.\n",
        "- Do NOT invent source numbers. Only use IDs that correspond to the provided source list.\n",
        "- There must be **no bullets or paragraphs without at least one citation** unless they are purely meta text (e.g., 'Bottom line:' heading).\n",
        "- If you summarize multiple sources in a single bullet, include all relevant numbers, e.g. `[1, 2, 4]`.\n",
        "- The documents are numbered in the order they appear below (1, 2, 3, etc.).\n",
        "- For example: 'Joint Cloud provides multi-level transport [1] with cost advantages [2].'\n",
        "- If referencing multiple sources: 'The system supports both IL5 and IL6 operations [1, 2].'\n",
        "\n",
        "CRITICAL: Every bullet point and paragraph that references information from the documents MUST end with at least one citation. This is mandatory.\n"
    ]
    
    # Load RAG files - if chat_id is provided, use it directly for efficiency
    all_files = []
    if chat_id:
        # Direct lookup using chat_id
        rag_files = load_rag_files(chat_id)
        for file_id in rag_file_ids:
            file = next((f for f in rag_files if f.get("id") == file_id), None)
            if file:
                all_files.append(file)
    else:
        # Fallback: search all RAG files (less efficient but works)
        for file_id in rag_file_ids:
            for rag_file_path in RAG_FILES_DIR.glob("*.json"):
                try:
                    with open(rag_file_path, 'r') as f:
                        files = json.load(f)
                        file = next((f for f in files if f.get("id") == file_id), None)
                        if file:
                            all_files.append(file)
                            break
                except:
                    continue
    
    # Load text for each file and build context
    # text_path is relative to project root (includes "uploads/" prefix)
    project_root = Path(__file__).parent.parent
    file_number = 1  # Track file number for citations (1-based to match RAG tray)
    for rag_file in all_files:
        if not rag_file.get("text_extracted") or not rag_file.get("text_path"):
            continue
        
        try:
            text_path = project_root / rag_file["text_path"]
            if text_path.exists():
                with open(text_path, 'r', encoding='utf-8') as f:
                    text_content = f.read()
                
                # Always include RAG files when available - they're explicitly uploaded for context
                # Use keyword relevance to prioritize chunks, but include all files
                user_keywords = set(user_message.lower().split())
                text_lower = text_content.lower()
                relevance_score = sum(1 for kw in user_keywords if kw in text_lower)
                
                chunks = chunk_text(text_content)
                # If there's relevance, take first 5 chunks, otherwise take first 3
                # This ensures files are always included but relevant ones get more context
                num_chunks = 5 if relevance_score > 0 else 3
                selected_chunks = chunks[:num_chunks]
                # Include file number in the source header so model knows which number to use for citations
                context_parts.append(f"\n----\nSource [{file_number}]: {rag_file['filename']}\n")
                context_parts.append('\n\n'.join(selected_chunks))
                context_parts.append("\n----\n")
                print(f"[RAG] Added context from {rag_file['filename']} ({len(selected_chunks)} chunks) as source [{file_number}]")
                file_number += 1
        except Exception as e:
            print(f"Error loading RAG file text: {e}")
            continue
    
    result = '\n'.join(context_parts)
    print(f"[RAG] Built context: {len(result)} characters, {len(all_files)} files processed")
    return result


def check_missing_citations(text: str) -> bool:
    """
    Check if RAG response has bullets/paragraphs missing citations.
    Returns True if citations are missing, False otherwise.
    """
    if not text:
        return False
    
    # Split into lines
    lines = text.split('\n')
    
    # Pattern to match citations like [1], [1, 3], [2, 4, 6]
    citation_pattern = re.compile(r'\[\s*\d+(?:\s*,\s*\d+)*\s*\]')
    
    # Track bullets and paragraphs
    in_bullet = False
    bullet_text = ""
    has_missing = False
    
    for line in lines:
        stripped = line.strip()
        
        # Skip empty lines and headers
        if not stripped or stripped.startswith('###'):
            if in_bullet and bullet_text:
                # Check if this bullet has citations
                if not citation_pattern.search(bullet_text):
                    # Check if it's not just meta text
                    if not any(meta in bullet_text.lower() for meta in ['bottom line', 'summary:', 'note:', 'important:']):
                        print(f"[RAG] Missing citation in bullet: {bullet_text[:100]}")
                        has_missing = True
                bullet_text = ""
                in_bullet = False
            continue
        
        # Check if this is a bullet point
        if stripped.startswith('- ') or stripped.startswith('* ') or stripped.startswith('• '):
            if in_bullet and bullet_text:
                # Check previous bullet
                if not citation_pattern.search(bullet_text):
                    if not any(meta in bullet_text.lower() for meta in ['bottom line', 'summary:', 'note:', 'important:']):
                        print(f"[RAG] Missing citation in bullet: {bullet_text[:100]}")
                        has_missing = True
            in_bullet = True
            bullet_text = stripped
        elif in_bullet:
            # Continuation of bullet (indented line)
            if stripped.startswith('  ') or stripped.startswith('\t'):
                bullet_text += " " + stripped
            else:
                # End of bullet, check it
                if not citation_pattern.search(bullet_text):
                    if not any(meta in bullet_text.lower() for meta in ['bottom line', 'summary:', 'note:', 'important:']):
                        print(f"[RAG] Missing citation in bullet: {bullet_text[:100]}")
                        has_missing = True
                in_bullet = False
                bullet_text = ""
        else:
            # Regular paragraph - check if it has citations
            if stripped and not stripped.startswith('###'):
                if not citation_pattern.search(stripped):
                    # Check if it's not just meta text
                    if not any(meta in stripped.lower() for meta in ['bottom line', 'summary:', 'note:', 'important:']):
                        # Only flag if it looks like content (not just a heading)
                        if len(stripped) > 20:  # Avoid flagging short headings
                            print(f"[RAG] Missing citation in paragraph: {stripped[:100]}")
                            has_missing = True
    
    # Check last bullet if still in progress
    if in_bullet and bullet_text:
        if not citation_pattern.search(bullet_text):
            if not any(meta in bullet_text.lower() for meta in ['bottom line', 'summary:', 'note:', 'important:']):
                print(f"[RAG] Missing citation in final bullet: {bullet_text[:100]}")
                has_missing = True
    
    return has_missing


async def fix_missing_citations(original_response: str, rag_context: str, model_display: str, provider: str) -> str:
    """
    Request a fix-up from the model to add missing citations.
    Returns the revised response.
    """
    fix_prompt = f"""Please revise the following answer so that every bullet and paragraph includes at least one inline source citation in `[n]` format at the end.

The original answer:
{original_response}

Please ensure:
- Every bullet point ends with at least one citation like [1] or [1, 3]
- Every paragraph that references information from sources ends with at least one citation
- Do not add citations to pure meta text (like "Bottom line:" headings)
- Keep all the original content and structure, just add the missing citations

Revised answer:"""
    
    messages = [
        {"role": "system", "content": "You are a helpful assistant that adds inline citations to text while preserving all original content."},
        {"role": "user", "content": fix_prompt}
    ]
    
    try:
        from chatdo.agents.main_agent import call_ai_router
        assistant_messages, _, _, _ = call_ai_router(messages, intent="general_chat")
        
        if assistant_messages and len(assistant_messages) > 0:
            fixed_response = assistant_messages[-1].get("content", original_response)
            print(f"[RAG] Fix-up completed, new length: {len(fixed_response)} chars")
            return fixed_response
        else:
            print(f"[RAG] Fix-up failed, returning original response")
            return original_response
    except Exception as e:
        print(f"[RAG] Error during fix-up: {e}")
        return original_response


@app.post("/api/rag/files", response_model=RagFile)
async def upload_rag_file(
    chat_id: str = Form(...),
    file: UploadFile = File(...)
):
    """
    Upload a file for RAG context in a specific chat
    """
    try:
        # Use handle_file_upload but store as RAG file
        # We'll use a special "rag" project_id
        upload_result = await handle_file_upload("rag", chat_id, file)
        
        # Create RagFile entry
        rag_file = RagFile(
            id=str(uuid.uuid4()),
            chat_id=chat_id,
            filename=upload_result.get("original_filename", upload_result.get("filename", "")),
            mime_type=upload_result.get("mime_type", "application/octet-stream"),
            size=upload_result.get("size", 0),
            created_at=datetime.now(timezone.utc).isoformat(),
            text_path=upload_result.get("text_path"),
            path=upload_result.get("path"),  # Store original file path for preview
            text_extracted=upload_result.get("text_extracted", False),
            error=upload_result.get("extraction_error")
        )
        
        # Save to RAG files list
        rag_files = load_rag_files(chat_id)
        rag_files.append(rag_file.dict())
        save_rag_files(chat_id, rag_files)
        
        return rag_file
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error uploading RAG file: {str(e)}")


@app.get("/api/rag/files")
async def list_rag_files(chat_id: str):
    """
    List all RAG files for a chat
    """
    rag_files = load_rag_files(chat_id)
    return rag_files


@app.get("/api/rag/find-original")
async def find_original_file(text_path: str, mime_type: str):
    """
    Find the original file path given a text_path.
    This is needed when old RAG files don't have the 'path' field.
    """
    from pathlib import Path
    
    # text_path format: uploads/rag/chat_id/uuid.txt
    # We need to find the original file in the same directory
    uploads_base = Path(__file__).parent.parent / "uploads"
    text_file_path = uploads_base / text_path.replace('uploads/', '')
    
    if not text_file_path.exists():
        raise HTTPException(status_code=404, detail="Text file not found")
    
    # Get the directory containing the text file
    file_dir = text_file_path.parent
    
    # Determine expected extension from mime type
    extension = ''
    if mime_type == 'application/pdf':
        extension = '.pdf'
    elif 'presentation' in mime_type or 'powerpoint' in mime_type:
        extension = '.pptx'
    elif 'spreadsheet' in mime_type or 'excel' in mime_type or 'sheet' in mime_type:
        extension = '.xlsx'
    elif 'word' in mime_type or 'wordprocessing' in mime_type:
        extension = '.docx'
    
    # Find all files in the directory with the expected extension (excluding .txt files)
    matching_files = [f for f in file_dir.iterdir() 
                      if f.is_file() 
                      and f.suffix == extension 
                      and not f.name.endswith('.txt')]
    
    if not matching_files:
        raise HTTPException(status_code=404, detail=f"Original file with extension {extension} not found in directory")
    
    # If multiple files, prefer the one that's not a .txt file
    # Usually there should be only one, but if there are multiple, take the first
    original_file = matching_files[0]
    
    # Return path relative to uploads base parent (to match the format used elsewhere)
    path = str(original_file.relative_to(uploads_base.parent))
    
    return {"path": path}


@app.delete("/api/rag/files/{file_id}")
async def delete_rag_file(chat_id: str, file_id: str):
    """
    Delete a RAG file
    """
    rag_files = load_rag_files(chat_id)
    original_count = len(rag_files)
    
    # Find the file to delete before filtering
    deleted_file = next((f for f in rag_files if f.get("id") == file_id), None)
    
    if not deleted_file:
        # File not found - might have been already deleted (e.g., in parallel delete)
        # Return success to avoid errors when clearing all files
        return {"success": True, "message": "File not found (may have been already deleted)"}
    
    # Filter out the deleted file
    rag_files = [f for f in rag_files if f.get("id") != file_id]
    
    # Delete the actual file and text file
    uploads_dir = Path(__file__).parent.parent / "uploads"
    if deleted_file.get("text_path"):
        text_path = uploads_dir / deleted_file["text_path"]
        if text_path.exists():
            try:
                text_path.unlink()
            except Exception as e:
                print(f"Warning: Failed to delete text file {text_path}: {e}")
    
    # Also try to delete the original file if path is available
    if deleted_file.get("path"):
        original_path = uploads_dir / deleted_file["path"]
        if original_path.exists():
            try:
                original_path.unlink()
            except Exception as e:
                print(f"Warning: Failed to delete original file {original_path}: {e}")
    
    save_rag_files(chat_id, rag_files)
    return {"success": True}


@app.get("/uploads/{file_path:path}")
async def serve_uploaded_file(file_path: str):
    """
    Serve uploaded files for preview
    """
    from fastapi.responses import FileResponse
    
    uploads_dir = Path(__file__).parent.parent / "uploads"
    file_full_path = uploads_dir / file_path
    
    # Security: ensure file is within uploads directory
    try:
        file_full_path.resolve().relative_to(uploads_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    
    if not file_full_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(file_full_path)


@app.get("/api/pptx-preview/{file_path:path}")
async def preview_pptx(file_path: str):
    """
    Convert PPTX to PDF for browser preview (same beautiful experience as PDFs!)
    Returns the PDF file for iframe display
    """
    from fastapi.responses import FileResponse
    import subprocess
    import tempfile
    
    uploads_dir = Path(__file__).parent.parent / "uploads"
    file_full_path = uploads_dir / file_path
    
    # Security check
    try:
        file_full_path.resolve().relative_to(uploads_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    
    if not file_full_path.exists() or not file_full_path.suffix.lower() in ['.pptx', '.ppt']:
        raise HTTPException(status_code=404, detail="File not found or not a PowerPoint file")
    
    # Check if we already have a converted PDF cached
    pdf_path = file_full_path.with_suffix('.pdf')
    if pdf_path.exists():
        return FileResponse(pdf_path, media_type='application/pdf')
    
    # Try to convert PPTX to PDF using LibreOffice (best quality)
    try:
        # Try multiple possible LibreOffice command paths
        libreoffice_commands = [
            '/opt/homebrew/bin/soffice',  # Homebrew on Apple Silicon (most common)
            '/usr/local/bin/soffice',  # Homebrew on Intel
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS app bundle
            'soffice',  # If in PATH
            'libreoffice',  # Alternative command name
        ]
        
        libreoffice_cmd = None
        for cmd in libreoffice_commands:
            try:
                # Check if command exists
                if cmd.startswith('/'):
                    # Absolute path - check if file exists
                    import os
                    if os.path.exists(cmd) and os.access(cmd, os.X_OK):
                        libreoffice_cmd = cmd
                        break
                else:
                    # Command name - use which
                    result = subprocess.run(
                        ['which', cmd],
                        capture_output=True,
                        timeout=2
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = result.stdout.decode().strip()
                        break
            except:
                continue
        
        if libreoffice_cmd:
            try:
                # Use LibreOffice headless to convert
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(pdf_path.parent),
                        str(file_full_path)
                    ],
                    capture_output=True,
                    timeout=60,
                    check=False
                )
                
                if result.returncode == 0 and pdf_path.exists():
                    return FileResponse(pdf_path, media_type='application/pdf')
                else:
                    # Conversion failed - log error but continue to fallback
                    error_msg = result.stderr.decode() if result.stderr else result.stdout.decode() if result.stdout else 'Unknown error'
                    print(f"LibreOffice conversion failed (code {result.returncode}): {error_msg}")
            except Exception as e:
                print(f"LibreOffice conversion error: {e}")
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
        # Fallback: Try python-pptx + reportlab (limited but works)
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas  # type: ignore
            from reportlab.lib.pagesizes import letter  # type: ignore
            from reportlab.lib.units import inch  # type: ignore
            
            prs = Presentation(file_full_path)
            
            # Create PDF
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            width, height = letter
            
            for slide_num, slide in enumerate(prs.slides, 1):
                if slide_num > 1:
                    c.showPage()
                
                y = height - 50
                # Add slide number
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, y, f"Slide {slide_num}")
                y -= 30
                
                # Extract and add title
                if slide.shapes.title and slide.shapes.title.text:
                    c.setFont("Helvetica-Bold", 14)
                    title_lines = slide.shapes.title.text.split('\n')
                    for line in title_lines[:3]:  # Limit to 3 lines
                        if y < 100:
                            break
                        c.drawString(50, y, line[:80])  # Limit width
                        y -= 20
                    y -= 10
                
                # Add text content
                c.setFont("Helvetica", 11)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                        text = shape.text.strip()
                        if text and y > 100:
                            # Wrap text
                            lines = text.split('\n')
                            for line in lines[:10]:  # Limit lines per slide
                                if y < 100:
                                    break
                                # Simple word wrap
                                words = line.split()
                                current_line = ""
                                for word in words:
                                    if len(current_line + word) < 80:
                                        current_line += word + " "
                                    else:
                                        if current_line:
                                            c.drawString(70, y, current_line.strip())
                                            y -= 15
                                        current_line = word + " "
                                if current_line and y > 100:
                                    c.drawString(70, y, current_line.strip())
                                    y -= 15
                                y -= 5
                
            c.save()
            
            if pdf_path.exists():
                return FileResponse(pdf_path, media_type='application/pdf')
            else:
                raise Exception("PDF creation failed")
                
        except ImportError:
            # No conversion libraries available
            raise HTTPException(
                status_code=501, 
                detail="PPTX to PDF conversion not available. Install LibreOffice or python libraries (reportlab, python-pptx) for preview support."
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error converting PPTX to PDF: {str(e)}")


@app.get("/api/xlsx-preview/{file_path:path}")
async def preview_xlsx(file_path: str):
    """
    Convert Excel (.xlsx, .xls) to PDF for browser preview (same beautiful experience as PDFs!)
    Returns the PDF file for iframe display
    """
    from fastapi.responses import FileResponse
    import subprocess
    
    uploads_dir = Path(__file__).parent.parent / "uploads"
    file_full_path = uploads_dir / file_path
    
    # Security check
    try:
        file_full_path.resolve().relative_to(uploads_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    
    if not file_full_path.exists() or not file_full_path.suffix.lower() in ['.xlsx', '.xls']:
        raise HTTPException(status_code=404, detail="File not found or not an Excel file")
    
    # Check if we already have a converted PDF cached
    pdf_path = file_full_path.with_suffix('.pdf')
    if pdf_path.exists():
        return FileResponse(pdf_path, media_type='application/pdf')
    
    # Try to convert Excel to PDF using LibreOffice (best quality)
    try:
        # Try multiple possible LibreOffice command paths
        libreoffice_commands = [
            '/opt/homebrew/bin/soffice',  # Homebrew on Apple Silicon (most common)
            '/usr/local/bin/soffice',  # Homebrew on Intel
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS app bundle
            'soffice',  # If in PATH
            'libreoffice',  # Alternative command name
        ]
        
        libreoffice_cmd = None
        for cmd in libreoffice_commands:
            try:
                # Check if command exists
                if cmd.startswith('/'):
                    # Absolute path - check if file exists
                    import os
                    if os.path.exists(cmd) and os.access(cmd, os.X_OK):
                        libreoffice_cmd = cmd
                        break
                else:
                    # Command name - use which
                    result = subprocess.run(
                        ['which', cmd],
                        capture_output=True,
                        timeout=2
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = result.stdout.decode().strip()
                        break
            except:
                continue
        
        if libreoffice_cmd:
            try:
                # Use LibreOffice headless to convert
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(pdf_path.parent),
                        str(file_full_path)
                    ],
                    capture_output=True,
                    timeout=60,
                    check=False
                )
                
                if result.returncode == 0 and pdf_path.exists():
                    return FileResponse(pdf_path, media_type='application/pdf')
                else:
                    # Conversion failed - log error but continue to fallback
                    error_msg = result.stderr.decode() if result.stderr else result.stdout.decode() if result.stdout else 'Unknown error'
                    print(f"LibreOffice Excel conversion failed (code {result.returncode}): {error_msg}")
            except Exception as e:
                print(f"LibreOffice Excel conversion error: {e}")
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
        pass
    
    # If all conversion attempts failed, return error
    raise HTTPException(status_code=500, detail="Excel to PDF conversion not available. Install LibreOffice for preview support.")


@app.get("/api/docx-preview/{file_path:path}")
async def preview_docx(file_path: str):
    """
    Convert Word (.docx, .doc) to PDF for browser preview (same beautiful experience as PDFs!)
    Returns the PDF file for iframe display
    """
    from fastapi.responses import FileResponse
    import subprocess
    
    uploads_dir = Path(__file__).parent.parent / "uploads"
    file_full_path = uploads_dir / file_path
    
    # Security check
    try:
        file_full_path.resolve().relative_to(uploads_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    
    if not file_full_path.exists() or not file_full_path.suffix.lower() in ['.docx', '.doc']:
        raise HTTPException(status_code=404, detail="File not found or not a Word document")
    
    # Check if we already have a converted PDF cached
    pdf_path = file_full_path.with_suffix('.pdf')
    if pdf_path.exists():
        return FileResponse(pdf_path, media_type='application/pdf')
    
    # Try to convert Word to PDF using LibreOffice (best quality)
    try:
        # Try multiple possible LibreOffice command paths
        libreoffice_commands = [
            '/opt/homebrew/bin/soffice',  # Homebrew on Apple Silicon (most common)
            '/usr/local/bin/soffice',  # Homebrew on Intel
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS app bundle
            'soffice',  # If in PATH
            'libreoffice',  # Alternative command name
        ]
        
        libreoffice_cmd = None
        for cmd in libreoffice_commands:
            try:
                # Check if command exists
                if cmd.startswith('/'):
                    # Absolute path - check if file exists
                    import os
                    if os.path.exists(cmd) and os.access(cmd, os.X_OK):
                        libreoffice_cmd = cmd
                        break
                else:
                    # Command name - use which
                    result = subprocess.run(
                        ['which', cmd],
                        capture_output=True,
                        timeout=2
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = result.stdout.decode().strip()
                        break
            except:
                continue
        
        if libreoffice_cmd:
            try:
                # Use LibreOffice headless to convert
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(pdf_path.parent),
                        str(file_full_path)
                    ],
                    capture_output=True,
                    timeout=60,
                    check=False
                )
                
                if result.returncode == 0 and pdf_path.exists():
                    return FileResponse(pdf_path, media_type='application/pdf')
                else:
                    # Conversion failed - log error but continue to fallback
                    error_msg = result.stderr.decode() if result.stderr else result.stdout.decode() if result.stdout else 'Unknown error'
                    print(f"LibreOffice Word conversion failed (code {result.returncode}): {error_msg}")
            except Exception as e:
                print(f"LibreOffice Word conversion error: {e}")
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
        pass
    
    # If all conversion attempts failed, return error
    raise HTTPException(status_code=500, detail="Word to PDF conversion not available. Install LibreOffice for preview support.")


@app.post("/api/url")
async def scrape_url_endpoint(
    project_id: str,
    conversation_id: str,
    url: str
):
    """
    Scrape URL content and save to uploads folder
    """
    try:
        result = await scrape_url(project_id, conversation_id, url)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/chats", response_model=List[Chat])
async def get_chats(
    project_id: Optional[str] = Query(None),
    include_trashed: bool = Query(False)
):
    """Get chats, optionally filtered by project and including trashed"""
    chats = load_chats()
    projects = load_projects()
    
    # Ensure General project exists
    general_project = next((p for p in projects if p.get("name") == "General"), None)
    if not general_project:
        general_project = {
            "id": str(uuid.uuid4()),
            "name": "General",
            "default_target": "general",
            "sort_index": 0,
            "created_at": now_iso(),
            "updated_at": now_iso()
        }
        projects.insert(0, general_project)
        save_projects(projects)
    
    # Assign orphaned chats (without project_id) to General
    needs_save = False
    for chat in chats:
        if not chat.get("project_id"):
            chat["project_id"] = general_project["id"]
            chat["updated_at"] = now_iso()
            needs_save = True
    
    if needs_save:
        save_chats(chats)
    
    # Filter by project if specified
    if project_id:
        chats = [c for c in chats if c.get("project_id") == project_id]
    
    # Filter trashed if not including them
    if not include_trashed:
        chats = get_active_chats(chats)
    
    return chats


@app.get("/api/chats/{chat_id}/messages")
async def get_chat_messages(chat_id: str, limit: Optional[int] = None):
    """Get messages for a specific chat conversation
    
    Args:
        chat_id: The chat ID
        limit: Optional limit on number of messages to return (for previews)
    """
    chats = load_chats()
    chat = next((c for c in chats if c.get("id") == chat_id), None)
    
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Get project to find target_name
    projects = load_projects()
    project = next((p for p in projects if p.get("id") == chat.get("project_id")), None)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    target_name = project.get("default_target", "general")
    thread_id = chat.get("thread_id") or chat_id  # Use chat_id as thread_id if thread_id not set
    
    print(f"[DIAG] get_chat_messages: chat_id={chat_id}, thread_id={thread_id}, target_name={target_name}")
    
    if not thread_id:
        print(f"[DIAG] WARNING: No thread_id for chat {chat_id}, returning empty messages")
        return {"messages": []}
    
    # Load messages from memory store
    history = load_thread_history(target_name, thread_id)
    print(f"[DIAG] get_chat_messages: Loaded {len(history)} messages from memory store")
    
    # Convert to frontend format
    messages = []
    print(f"[DIAG] Loading messages for chat {chat_id}: history has {len(history)} messages")
    for idx, msg in enumerate(history):
        # Skip system messages for display
        if msg.get("role") == "system":
            print(f"[DIAG] Skipping system message at index {idx}")
            continue
        
        # Preserve structured message types if they exist
        message_obj = {
            "role": msg.get("role"),
            "content": msg.get("content", "")
        }
        
        # If message has structured type/data, preserve it
        if msg.get("type"):
            message_obj["type"] = msg.get("type")
        if msg.get("data"):
            message_obj["data"] = msg.get("data")
        if msg.get("model"):
            message_obj["model"] = msg.get("model")
        if msg.get("provider"):
            message_obj["provider"] = msg.get("provider")
        if msg.get("sources"):
            message_obj["sources"] = msg.get("sources")
        
        print(f"[DIAG] Message {idx}: role={msg.get('role')}, content_length={len(msg.get('content', ''))}, type={msg.get('type', 'none')}, content_preview={msg.get('content', '')[:50]}...")
        messages.append(message_obj)
    
    print(f"[DIAG] Returning {len(messages)} messages to frontend")
    
    # If limit is specified, return only the last N messages (for previews)
    if limit and limit > 0:
        messages = messages[-limit:]
    
    return {"messages": messages}


@app.get("/api/chats/{chat_id}/sources")
async def get_chat_sources(chat_id: str):
    """Get sources for a specific chat conversation"""
    chats = load_chats()
    chat = next((c for c in chats if c.get("id") == chat_id), None)
    
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Get project to find target_name
    projects = load_projects()
    project = next((p for p in projects if p.get("id") == chat.get("project_id")), None)
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    target_name = project.get("default_target", "general")
    thread_id = chat.get("thread_id")
    
    if not thread_id:
        return {"sources": []}
    
    # Load sources from memory store
    sources = load_thread_sources(target_name, thread_id)
    
    return {"sources": sources}


@app.get("/favicon.ico")
async def favicon():
    """Handle favicon requests to prevent 410 errors"""
    # Return 204 No Content to indicate no favicon is available
    return JSONResponse(status_code=204, content=None)


@app.patch("/api/chats/{chat_id}", response_model=Chat)
async def update_chat(chat_id: str, chat_data: ChatUpdate):
    """Update a chat (e.g., rename)"""
    chats = load_chats()
    
    # Find chat by ID
    chat_index = None
    for i, c in enumerate(chats):
        if c.get("id") == chat_id:
            chat_index = i
            break
    
    if chat_index is None:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Update chat fields
    if chat_data.title is not None:
        chats[chat_index]["title"] = chat_data.title.strip()
        chats[chat_index]["updated_at"] = now_iso()
    
    save_chats(chats)
    return chats[chat_index]


@app.delete("/api/chats/{chat_id}", response_model=Chat)
async def soft_delete_chat(chat_id: str):
    """Soft delete a chat (move to trash)"""
    chats = load_chats()
    
    # Find chat by ID
    chat_index = None
    for i, c in enumerate(chats):
        if c.get("id") == chat_id:
            chat_index = i
            break
    
    if chat_index is None:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Soft delete
    chats[chat_index]["trashed"] = True
    chats[chat_index]["trashed_at"] = now_iso()
    chats[chat_index]["updated_at"] = now_iso()
    
    save_chats(chats)
    return chats[chat_index]


@app.post("/api/chats/{chat_id}/restore", response_model=Chat)
async def restore_chat(chat_id: str):
    """Restore a chat from trash"""
    chats = load_chats()
    
    # Find chat by ID
    chat_index = None
    for i, c in enumerate(chats):
        if c.get("id") == chat_id:
            chat_index = i
            break
    
    if chat_index is None:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Restore
    chats[chat_index]["trashed"] = False
    chats[chat_index]["trashed_at"] = None
    chats[chat_index]["updated_at"] = now_iso()
    
    save_chats(chats)
    return chats[chat_index]


@app.post("/api/chats/{chat_id}/purge")
async def purge_chat(chat_id: str):
    """Permanently delete a chat and its thread history"""
    chats = load_chats()
    
    # Find chat by ID
    chat = None
    for c in chats:
        if c.get("id") == chat_id:
            chat = c
            break
    
    if chat is None:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Delete thread history
    project_id = chat.get("project_id")
    thread_id = chat.get("thread_id")
    
    if project_id and thread_id:
        projects = load_projects()
        project = next((p for p in projects if p.get("id") == project_id), None)
        if project:
            target_name = project.get("default_target", "general")
            try:
                delete_thread_history(target_name, thread_id)
            except Exception as e:
                print(f"Warning: Failed to delete thread history for {thread_id}: {e}")
    
    # Remove from chats list
    chats = [c for c in chats if c.get("id") != chat_id]
    save_chats(chats)
    
    return {"success": True}


@app.post("/api/chats/purge_trashed")
async def purge_trashed_endpoint():
    """Manually trigger purge of old trashed chats"""
    purged_count = purge_trashed_chats()
    return {"success": True, "purged_count": purged_count}


@app.get("/api/projects/{project_id}/memory-sources")
async def get_project_memory_sources(project_id: str):
    """Get memory sources for a project"""
    from server.services import projects_config
    
    project = projects_config.get_project(project_id)
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Always return a list
    sources = project.get("memory_sources") or []
    return {"project_id": project_id, "memory_sources": sources}


class UpdateMemorySourcesRequest(BaseModel):
    memory_sources: List[str]


class AddMemorySourceRequest(BaseModel):
    root_path: str
    display_name: Optional[str] = None
    project_id: Optional[str] = "scratch"


@app.put("/api/projects/{project_id}/memory-sources")
async def update_project_memory_sources(project_id: str, body: UpdateMemorySourcesRequest):
    """Update memory sources for a project"""
    from server.services import projects_config
    from memory_service.config import sync_yaml_from_projects
    
    try:
        project = projects_config.update_project_memory_sources(
            project_id, body.memory_sources
        )
        
        # Auto-sync YAML file to reflect the new connections
        sync_yaml_from_projects()
        
        return {
            "project_id": project_id,
            "memory_sources": project.get("memory_sources") or [],
        }
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))


@app.post("/api/memory/sources")
async def api_add_memory_source(req: AddMemorySourceRequest):
    """Proxy: UI → server → memory_service.add_source"""
    from server.services import memory_service_client
    
    try:
        client = memory_service_client.get_memory_client()
        result = client.add_memory_source(
            root_path=req.root_path,
            display_name=req.display_name,
            project_id=req.project_id,
        )
        return result
    except Exception as e:
        detail = str(e)
        if "Memory Service is not available" in detail:
            raise HTTPException(status_code=503, detail=detail)
        raise HTTPException(status_code=400, detail=detail)


@app.delete("/api/memory/sources/{source_id}")
async def api_delete_memory_source(source_id: str):
    """Proxy: UI → server → memory_service.delete_source"""
    from server.services import memory_service_client
    
    try:
        client = memory_service_client.get_memory_client()
        result = client.delete_memory_source(source_id)
        return result
    except Exception as e:
        detail = str(e)
        if "Memory Service is not available" in detail:
            raise HTTPException(status_code=503, detail=detail)
        raise HTTPException(status_code=400, detail=detail)


@app.post("/api/chats/purge_all_trashed")
async def purge_all_trashed():
    """Permanently delete ALL trashed chats (regardless of age)"""
    chats = load_chats()
    projects = load_projects()
    
    # Find all trashed chats
    trashed_chat_ids = []
    for chat in chats:
        if chat.get("trashed", False):
            trashed_chat_ids.append(chat.get("id"))
            # Get project to find target_name for thread deletion
            project = next((p for p in projects if p.get("id") == chat.get("project_id")), None)
            target_name = project.get("default_target", "general") if project else "general"
            thread_id = chat.get("thread_id")
            
            # Delete thread history if it exists
            if thread_id:
                try:
                    delete_thread_history(target_name, thread_id)
                except Exception as e:
                    print(f"Warning: Failed to delete thread history for {thread_id}: {e}")
    
    # Remove all trashed chats from the list
    chats = [c for c in chats if not c.get("trashed", False)]
    save_chats(chats)
    
    return {"success": True, "purged_count": len(trashed_chat_ids)}


# Impact endpoints
class ImpactCreate(BaseModel):
    title: str
    date: Optional[date] = None
    context: Optional[str] = None
    actions: str
    impact: Optional[str] = None
    metrics: Optional[str] = None
    tags: list[str] = []
    notes: Optional[str] = None


class ImpactUpdate(BaseModel):
    title: Optional[str] = None
    date: Optional[date] = None
    context: Optional[str] = None
    actions: Optional[str] = None
    impact: Optional[str] = None
    metrics: Optional[str] = None
    tags: Optional[list[str]] = None
    notes: Optional[str] = None


@app.get("/api/impacts/", response_model=List[ImpactEntry])
def get_impacts():
    return impact_store.list_impacts()


@app.post("/api/impacts/", response_model=ImpactEntry)
def create_impact(entry: ImpactCreate):
    return impact_store.create_impact(
        ImpactEntry(**entry.model_dump())
    )


@app.put("/api/impacts/{impact_id}", response_model=ImpactEntry)
def update_impact(impact_id: str, patch: ImpactUpdate):
    updated = impact_store.update_impact(impact_id, patch.model_dump(exclude_unset=True))
    if not updated:
        raise HTTPException(status_code=404, detail="Impact not found")
    return updated


@app.delete("/api/impacts/{impact_id}")
def delete_impact(impact_id: str):
    ok = impact_store.delete_impact(impact_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Impact not found")
    return {"ok": True}


# Impact template endpoints
class ImpactTemplateResponse(BaseModel):
    id: str
    created_at: datetime
    name: str
    description: Optional[str]
    tags: list[str]
    file_name: str
    mime_type: Optional[str]


@app.get("/api/impact-templates/", response_model=List[ImpactTemplateResponse])
def list_impact_templates():
    templates = impact_templates_store.list_templates()
    return [
        ImpactTemplateResponse(
            id=t.id,
            created_at=t.created_at,
            name=t.name,
            description=t.description,
            tags=t.tags,
            file_name=t.file_name,
            mime_type=t.mime_type,
        )
        for t in templates
    ]


@app.post("/api/impact-templates/upload", response_model=ImpactTemplateResponse)
async def upload_impact_template(
    file: UploadFile = File(...),
    name: str = Form(...),
    description: Optional[str] = Form(None),
    tags: Optional[str] = Form(None),  # comma-separated string from UI
):
    impact_templates_store.ensure_paths()  # ensure dirs
    safe_id = uuid4().hex
    # keep extension
    from pathlib import Path as P
    ext = P(file.filename).suffix or ""
    stored_name = f"{safe_id}{ext}"
    dest = impact_templates_store.TEMPLATES_DIR / stored_name
    content = await file.read()
    dest.write_bytes(content)
    
    tag_list: list[str] = []
    if tags:
        tag_list = [t.strip() for t in tags.split(",") if t.strip()]
    
    tpl = ImpactTemplate(
        id=safe_id,
        name=name,
        description=description,
        tags=tag_list,
        file_name=file.filename,
        stored_name=stored_name,
        mime_type=file.content_type,
    )
    tpl = impact_templates_store.add_template(tpl)
    return ImpactTemplateResponse(
        id=tpl.id,
        created_at=tpl.created_at,
        name=tpl.name,
        description=tpl.description,
        tags=tpl.tags,
        file_name=tpl.file_name,
        mime_type=tpl.mime_type,
    )


@app.delete("/api/impact-templates/{template_id}")
def delete_impact_template(template_id: str):
    # For now, just remove from JSON. We can later also delete the file.
    ok = impact_templates_store.delete_template(template_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Template not found")
    return {"ok": True}


@app.get("/api/impact-templates/{template_id}/file")
def download_impact_template_file(template_id: str):
    path = get_template_file_path(template_id)
    if not path or not path.exists():
        raise HTTPException(status_code=404, detail="Template file not found")
    return FileResponse(path)


# Template Autofill Engine endpoints
class TemplateUploadResponse(BaseModel):
    template_id: str
    filename: str
    fields: List[Dict[str, Any]]


@app.post("/api/templates/upload", response_model=TemplateUploadResponse)
async def upload_template(file: UploadFile = File(...)):
    """Upload a template, extract text, and analyze for fillable fields."""
    try:
        # Save template
        file_content = await file.read()
        metadata = template_store.save_template(file_content, file.filename)
        
        # For PDFs, try to extract AcroForm fields first
        fields = []
        pages = None
        if metadata.ext.lower() == '.pdf':
            try:
                # Extract PDF form fields
                pdf_fields = extract_pdf_fields.extract_pdf_fields(metadata.path)
                if pdf_fields:
                    fields = pdf_fields
                    
                    # Count pages
                    try:
                        import PyPDF2
                        with open(metadata.path, 'rb') as f:
                            pdf_reader = PyPDF2.PdfReader(f)
                            pages = len(pdf_reader.pages)
                    except:
                        pass
            except Exception as e:
                logger.warning(f"PDF field extraction failed, falling back to text analysis: {e}")
        
        # If no PDF fields found, use text analysis
        if not fields:
            # Extract text using Unstructured
            template_text = extract_template_text.extract_template_text(metadata.path)
            if not template_text:
                raise HTTPException(status_code=400, detail="Could not extract text from template")
            
            # Analyze template to identify fields (synchronous function)
            fields = analyze_template.analyze_template(template_text)
        
        # Update template metadata with fields and pages
        template_store.update_template_fields(metadata.template_id, fields)
        if pages:
            # Update pages count
            from server.services.template_engine.template_store import _load_metadata, _save_metadata
            all_metadata = _load_metadata()
            if metadata.template_id in all_metadata:
                all_metadata[metadata.template_id]["pages"] = pages
                _save_metadata(all_metadata)
        
        return TemplateUploadResponse(
            template_id=metadata.template_id,
            filename=metadata.filename,
            fields=fields
        )
    except Exception as e:
        logger.error(f"Template upload failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Template upload failed: {str(e)}")


@app.get("/api/templates")
def list_templates():
    """List all uploaded templates with their fields."""
    templates = template_store.list_templates()
    return [
        {
            "template_id": t.template_id,
            "filename": t.filename,
            "ext": t.ext,
            "fields": t.fields,
            "created_at": t.created_at.isoformat(),
            "pages": t.pages,
        }
        for t in templates
    ]


@app.get("/api/templates/{template_id}")
def get_template(template_id: str):
    """Get a specific template by ID."""
    template = template_store.get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    return {
        "template_id": template.template_id,
        "filename": template.filename,
        "ext": template.ext,
        "fields": template.fields,
        "created_at": template.created_at.isoformat(),
        "pages": template.pages,
    }


@app.delete("/api/templates/{template_id}")
def delete_template(template_id: str):
    """Delete a template and its file."""
    template = template_store.get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # Delete the template directory (includes file and autofill outputs)
    import shutil
    template_dir = template.path.parent
    if template_dir.exists():
        shutil.rmtree(template_dir)
    
    # Remove from metadata - need to access internal functions
    from server.services.template_engine.template_store import _load_metadata, _save_metadata
    all_metadata = _load_metadata()
    if template_id in all_metadata:
        del all_metadata[template_id]
        _save_metadata(all_metadata)
    
    return {"ok": True}


class AutofillRequest(BaseModel):
    field_assignments: Dict[str, str]  # field_id -> impact_id


@app.post("/api/templates/{template_id}/autofill")
async def autofill_template_endpoint(template_id: str, request: AutofillRequest):
    """Autofill a template using assigned impacts."""
    try:
        # Get template
        template = template_store.get_template(template_id)
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")
        
        # Extract template text
        template_text = extract_template_text.extract_template_text(template.path)
        if not template_text:
            raise HTTPException(status_code=400, detail="Could not extract text from template")
        
        # Get all impacts
        impacts = impact_store.list_impacts()
        impacts_dict = [impact.model_dump() for impact in impacts]
        
        # Autofill template
        output_text = autofill_template.autofill_template(
            template_text=template_text,
            fields=template.fields,
            field_assignments=request.field_assignments,
            impacts=impacts_dict,
        )
        
        # Save autofilled output
        output_path = template_store.save_autofilled_output(template_id, output_text)
        
        return {
            "template_id": template_id,
            "output_path": str(output_path),
            "output_text": output_text,
        }
    except Exception as e:
        logger.error(f"Template autofill failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Template autofill failed: {str(e)}")


@app.get("/api/templates/{template_id}/autofill/latest")
def get_latest_autofill(template_id: str):
    """Get the latest autofilled output for a template."""
    latest_path = template_store.get_latest_autofill(template_id)
    if not latest_path or not latest_path.exists():
        raise HTTPException(status_code=404, detail="No autofilled output found")
    
    output_text = latest_path.read_text(encoding="utf-8")
    return {
        "template_id": template_id,
        "output_path": str(latest_path),
        "output_text": output_text,
        "created_at": datetime.fromtimestamp(latest_path.stat().st_mtime).isoformat(),
    }


@app.websocket("/api/chat/stream")
async def websocket_chat(websocket: WebSocket):
    """
    WebSocket endpoint for streaming chat responses
    """
    await websocket_endpoint(websocket)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

