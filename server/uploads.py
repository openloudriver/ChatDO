"""
File upload handling for ChatDO
Saves files and extracts text from PDFs, Word docs, images, etc.
"""
from fastapi import UploadFile
from pathlib import Path
import uuid
import aiofiles
import mimetypes


UPLOADS_BASE = Path(__file__).parent.parent / "uploads"


async def handle_file_upload(project_id: str, conversation_id: str, file: UploadFile) -> dict:
    """
    Save uploaded file and extract text if applicable.
    Returns metadata about the saved file.
    """
    # Create upload directory
    upload_dir = UPLOADS_BASE / project_id / conversation_id
    upload_dir.mkdir(parents=True, exist_ok=True)
    
    # Generate unique filename
    file_ext = Path(file.filename).suffix if file.filename else ""
    unique_filename = f"{uuid.uuid4()}{file_ext}"
    file_path = upload_dir / unique_filename
    
    # Save file
    async with aiofiles.open(file_path, 'wb') as f:
        content = await file.read()
        await f.write(content)
    
    # Determine MIME type
    mime_type, _ = mimetypes.guess_type(file.filename or "")
    
    result = {
        "filename": unique_filename,
        "original_filename": file.filename,
        "path": str(file_path.relative_to(UPLOADS_BASE.parent)),
        "size": len(content),
        "mime_type": mime_type,
        "text_extracted": False,
        "extracted_text": None
    }
    
    # Extract text based on file type
    if mime_type:
        if mime_type.startswith("text/"):
            # Plain text file
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                text_content = await f.read()
            text_path = upload_dir / f"{uuid.uuid4()}.txt"
            async with aiofiles.open(text_path, 'w', encoding='utf-8') as f:
                await f.write(text_content)
            result["text_extracted"] = True
            result["text_path"] = str(text_path.relative_to(UPLOADS_BASE.parent))
            result["extracted_text"] = text_content
        
        elif mime_type == "application/pdf":
            # PDF extraction (requires PyPDF2 or pdfplumber)
            try:
                text_content = await extract_pdf_text(file_path)
                if text_content and text_content.strip():
                    text_path = upload_dir / f"{uuid.uuid4()}.txt"
                    async with aiofiles.open(text_path, 'w', encoding='utf-8') as f:
                        await f.write(text_content)
                    result["text_extracted"] = True
                    result["text_path"] = str(text_path.relative_to(UPLOADS_BASE.parent))
                    result["extracted_text"] = text_content
                else:
                    result["extraction_error"] = "No text could be extracted from PDF (may be image-based/scanned)"
                    print(f"PDF extraction returned empty text for {file.filename}")
            except ImportError as e:
                result["extraction_error"] = f"PDF extraction library not installed: {str(e)}"
                print(f"PDF extraction failed: {e}")
            except Exception as e:
                result["extraction_error"] = str(e)
                print(f"PDF extraction error: {e}")
        
        elif mime_type in [
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ]:
            # Word document extraction
            try:
                text_content = await extract_word_text(file_path)
                if text_content and text_content.strip():
                    text_path = upload_dir / f"{uuid.uuid4()}.txt"
                    async with aiofiles.open(text_path, 'w', encoding='utf-8') as f:
                        await f.write(text_content)
                    result["text_extracted"] = True
                    result["text_path"] = str(text_path.relative_to(UPLOADS_BASE.parent))
                    result["extracted_text"] = text_content
                else:
                    result["extraction_error"] = "No text could be extracted from Word document"
                    print(f"Word document extraction returned empty text for {file.filename}")
            except ImportError as e:
                result["extraction_error"] = f"Word document extraction library not installed: {str(e)}"
                print(f"Word document extraction failed: {e}")
            except Exception as e:
                result["extraction_error"] = str(e)
                print(f"Word document extraction error: {e}")
        
        elif mime_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            # PowerPoint extraction
            try:
                text_content = await extract_pptx_text(file_path)
                if text_content:
                    text_path = upload_dir / f"{uuid.uuid4()}.txt"
                    async with aiofiles.open(text_path, 'w', encoding='utf-8') as f:
                        await f.write(text_content)
                    result["text_extracted"] = True
                    result["text_path"] = str(text_path.relative_to(UPLOADS_BASE.parent))
                    result["extracted_text"] = text_content
            except Exception as e:
                result["extraction_error"] = str(e)
        
        elif mime_type.startswith("image/"):
            # Image OCR (would require pytesseract or similar)
            # For now, just save the image
            result["note"] = "Image OCR not yet implemented"
    
    return result


async def extract_pdf_text(pdf_path: Path) -> str:
    """Extract text from PDF file"""
    try:
        import PyPDF2
        text_parts = []
        with open(pdf_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        result = "\n\n".join(text_parts)
        if not result:
            print(f"Warning: PyPDF2 extracted no text from {pdf_path}")
        return result
    except ImportError:
        print("Warning: PyPDF2 not installed, trying pdfplumber...")
        # Fallback: try pdfplumber
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text or "")
            result = "\n\n".join(text_parts)
            if not result:
                print(f"Warning: pdfplumber extracted no text from {pdf_path}")
            return result
        except ImportError:
            print("Error: Neither PyPDF2 nor pdfplumber is installed. Please install one: pip install PyPDF2 or pip install pdfplumber")
            raise ImportError("No PDF extraction library available. Install PyPDF2 or pdfplumber.")
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        raise


async def extract_word_text(doc_path: Path) -> str:
    """Extract text from Word document"""
    try:
        from docx import Document
        doc = Document(doc_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result = "\n\n".join(paragraphs)
        if not result:
            print(f"Warning: python-docx extracted no text from {doc_path}")
        return result
    except ImportError:
        print("Error: python-docx not installed. Please install: pip install python-docx")
        raise ImportError("python-docx not available. Install with: pip install python-docx")
    except Exception as e:
        print(f"Error extracting Word document text: {e}")
        raise


async def extract_pptx_text(pptx_path: Path) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            
            # Extract title
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    # Check if it's a text box or placeholder
                    if shape.text.strip():
                        slide_text.append(shape.text.strip())
            
            # Extract speaker notes if available
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text.append(f"Notes: {notes_text}")
            
            if len(slide_text) > 1:  # More than just "Slide X:"
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except ImportError:
        # Try alternative: python-pptx might not be installed
        try:
            import zipfile
            # PPTX is a ZIP file - extract text from XML
            text_parts = []
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                # Get slide files
                slide_files = [f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                for slide_file in sorted(slide_files):
                    slide_num = slide_file.split('slide')[1].split('.')[0]
                    content = zip_ref.read(slide_file).decode('utf-8')
                    # Simple text extraction from XML (basic)
                    import re
                    text_matches = re.findall(r'<a:t[^>]*>([^<]+)</a:t>', content)
                    if text_matches:
                        text_parts.append(f"Slide {slide_num}:\n" + "\n".join(text_matches))
            return "\n\n".join(text_parts) if text_parts else ""
        except Exception:
            return ""  # python-pptx not available and ZIP extraction failed

